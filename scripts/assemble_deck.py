#!/usr/bin/env python3
"""Assemble a palette-controlled image deck from a validated deck plan."""

from __future__ import annotations

import argparse
import copy
import colorsys
import hashlib
import io
import json
import re
import tempfile
import zipfile
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from deck_utils import (
    configured_minimum_size,
    load_plan,
    resolve_path,
    resolved_title_render_mode,
    validate_image_geometry,
)
from validate_plan import validate_plan


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_TEMPLATE = ROOT / "assets" / "parallel-digital-standard-template.pptx"
HEX_IN_TEXT = re.compile(r"(?i)#([0-9a-f]{6})")


def replace_text_preserve_style(shape: Any, new_text: str) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    paragraphs = shape.text_frame.paragraphs
    first_run = next(
        (run for paragraph in paragraphs for run in paragraph.runs),
        None,
    )
    if first_run is None:
        shape.text = new_text
        return
    used = False
    for paragraph in paragraphs:
        for run in paragraph.runs:
            if not used:
                run.text = new_text
                used = True
            else:
                run.text = ""


def replace_matching_text(slide: Any, replacements: dict[str, str]) -> None:
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text.strip()
        name = getattr(shape, "name", "")
        for key, value in replacements.items():
            if text == key or name == key:
                replace_text_preserve_style(shape, value)


def duplicate_slide(prs: Presentation, source_slide: Any) -> Any:
    new_slide = prs.slides.add_slide(source_slide.slide_layout)
    for shape in list(new_slide.shapes):
        new_slide.shapes._spTree.remove(shape.element)
    for shape in source_slide.shapes:
        new_el = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
    return new_slide


def delete_slide_by_id(prs: Presentation, slide_id: int) -> None:
    for index, slide in enumerate(prs.slides):
        if slide.slide_id == slide_id:
            slide_id_element = prs.slides._sldIdLst[index]
            prs.part.drop_rel(slide_id_element.rId)
            prs.slides._sldIdLst.remove(slide_id_element)
            return


def reorder_slides(prs: Presentation, ordered_slide_ids: list[int]) -> None:
    slide_id_list = prs.slides._sldIdLst
    elements = {int(element.get("id")): element for element in list(slide_id_list)}
    for element in list(slide_id_list):
        slide_id_list.remove(element)
    for slide_id in ordered_slide_ids:
        element = elements.get(slide_id)
        if element is not None:
            slide_id_list.append(element)


def add_full_bleed_image_slide(prs: Presentation, image_path: Path) -> Any:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for shape in list(slide.shapes):
        slide.shapes._spTree.remove(shape.element)
    slide.shapes.add_picture(
        str(image_path),
        0,
        0,
        width=prs.slide_width,
        height=prs.slide_height,
    )
    return slide


def add_native_content_title(
    slide: Any,
    plan: dict[str, Any],
    slide_plan: dict[str, Any],
) -> None:
    title = str(slide_plan.get("title", "")).strip()
    if not title:
        return
    palette = plan["palette"]
    title_style = plan.get("native_title_style", {})
    font_name = str(title_style.get("font_name", "Microsoft YaHei"))
    font_size = float(title_style.get("font_size", 24))
    if len(title) > 48:
        font_size = min(font_size, 18)
    elif len(title) > 34:
        font_size = min(font_size, 20)

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.20),
        Inches(0.20),
        Inches(0.05),
        Inches(0.48),
    )
    accent.name = "Native Title Accent"
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor.from_string(
        palette["accent"].lstrip("#")
    )
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.34),
        Inches(0.13),
        Inches(12.55),
        Inches(0.66),
    )
    title_box.name = "Native Content Title"
    text_frame = title_box.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = title
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string(palette["text"].lstrip("#"))


def update_cover(slide: Any, plan: dict[str, Any]) -> None:
    replacements: dict[str, str] = {}
    if plan.get("cover_title"):
        replacements["平行数字PPT模板"] = plan["cover_title"]
        replacements["文本框 8"] = plan["cover_title"]
    if plan.get("footer"):
        replacements["平行数字  交叉现实"] = plan["footer"]
    replace_matching_text(slide, replacements)


def update_catalogue(
    slide: Any,
    sections: list[dict[str, Any]],
    footer: str | None,
) -> None:
    replacements: dict[str, str] = {}
    slots = ["第一节", "第二节", "第三节", "第四节"]
    for index, section in enumerate(sections):
        replacements[slots[index]] = (
            section.get("subtitle") or section.get("main_title") or slots[index]
        )
    if footer:
        replacements["平行数字  交叉现实"] = footer
    replace_matching_text(slide, replacements)


def color_tuple(value: str) -> tuple[int, int, int]:
    value = value.lstrip("#")
    return tuple(int(value[index : index + 2], 16) for index in (0, 2, 4))


def recolor_pixel(
    rgba: tuple[int, int, int, int],
    primary: tuple[int, int, int],
    accent: tuple[int, int, int],
) -> tuple[int, int, int, int]:
    red, green, blue, alpha = rgba
    if alpha == 0:
        return rgba
    hue, lightness, saturation = colorsys.rgb_to_hls(
        red / 255,
        green / 255,
        blue / 255,
    )
    if saturation < 0.16:
        return rgba
    if 0.44 <= hue <= 0.56:
        target = accent
    elif 0.53 <= hue <= 0.75:
        target = primary
    else:
        return rgba
    target_h, _, target_s = colorsys.rgb_to_hls(*(channel / 255 for channel in target))
    new_red, new_green, new_blue = colorsys.hls_to_rgb(
        target_h,
        lightness,
        max(target_s, saturation * 0.7),
    )
    return (
        round(new_red * 255),
        round(new_green * 255),
        round(new_blue * 255),
        alpha,
    )


def recolor_png(data: bytes, palette: dict[str, str]) -> bytes:
    try:
        from PIL import Image
    except ImportError as exc:
        raise RuntimeError("Pillow is required to recolor template raster assets") from exc
    primary = color_tuple(palette["primary"])
    accent = color_tuple(palette["accent"])
    with Image.open(io.BytesIO(data)) as image:
        rgba = image.convert("RGBA")
        pixels = (
            rgba.get_flattened_data()
            if hasattr(rgba, "get_flattened_data")
            else rgba.getdata()
        )
        cache: dict[
            tuple[int, int, int, int],
            tuple[int, int, int, int],
        ] = {}
        recolored = []
        for pixel in pixels:
            if pixel not in cache:
                cache[pixel] = recolor_pixel(pixel, primary, accent)
            recolored.append(cache[pixel])
        rgba.putdata(recolored)
        output = io.BytesIO()
        rgba.save(output, format="PNG", optimize=True)
        return output.getvalue()


def recolor_svg(data: bytes, palette: dict[str, str]) -> bytes:
    text = data.decode("utf-8")

    def replace(match: re.Match[str]) -> str:
        original = match.group(1)
        red, green, blue = color_tuple(original)
        hue, _, saturation = colorsys.rgb_to_hls(
            red / 255,
            green / 255,
            blue / 255,
        )
        if saturation < 0.16:
            return match.group(0)
        if 0.44 <= hue <= 0.56:
            return f"#{palette['accent'].lstrip('#')}"
        if 0.53 <= hue <= 0.75:
            return f"#{palette['primary'].lstrip('#')}"
        return match.group(0)

    return HEX_IN_TEXT.sub(replace, text).encode("utf-8")


def recolor_theme_xml(data: bytes, palette: dict[str, str]) -> bytes:
    text = data.decode("utf-8")
    replacements = {
        "dk2": palette["secondary"],
        "accent1": palette["primary"],
        "accent2": palette["accent"],
        "accent5": palette.get("muted", palette["secondary"]),
    }
    for name, color in replacements.items():
        pattern = rf'(<a:{name}><a:srgbClr val=")[0-9A-Fa-f]{{6}}'
        text = re.sub(pattern, rf"\g<1>{color.lstrip('#').upper()}", text)
    return text.encode("utf-8")


def prepare_template_with_palette(
    source: Path,
    target: Path,
    palette: dict[str, str],
) -> None:
    with zipfile.ZipFile(source, "r") as input_zip:
        with zipfile.ZipFile(target, "w", zipfile.ZIP_DEFLATED) as output_zip:
            for item in input_zip.infolist():
                data = input_zip.read(item.filename)
                lower = item.filename.lower()
                if lower.startswith("ppt/media/") and lower.endswith(".png"):
                    data = recolor_png(data, palette)
                elif lower.startswith("ppt/media/") and lower.endswith(".svg"):
                    data = recolor_svg(data, palette)
                elif lower.startswith("ppt/theme/") and lower.endswith(".xml"):
                    data = recolor_theme_xml(data, palette)
                output_zip.writestr(item, data)


def replace_rgb_color(color: Any, mapping: dict[str, str]) -> int:
    try:
        if color.type == MSO_COLOR_TYPE.RGB:
            source = str(color.rgb).upper()
            if source in mapping:
                color.rgb = RGBColor.from_string(mapping[source].lstrip("#"))
                return 1
    except (AttributeError, ValueError):
        pass
    return 0


def recolor_shape(shape: Any, mapping: dict[str, str]) -> int:
    changed = 0
    try:
        if shape.fill.type == MSO_FILL.SOLID:
            changed += replace_rgb_color(shape.fill.fore_color, mapping)
    except (AttributeError, ValueError):
        pass
    try:
        changed += replace_rgb_color(shape.line.color, mapping)
    except (AttributeError, ValueError):
        pass
    if getattr(shape, "has_text_frame", False):
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                changed += replace_rgb_color(run.font.color, mapping)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.fill.type == MSO_FILL.SOLID:
                    changed += replace_rgb_color(cell.fill.fore_color, mapping)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        changed += replace_rgb_color(run.font.color, mapping)
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            changed += recolor_shape(child, mapping)
    return changed


def recolor_native_shapes(prs: Presentation, palette: dict[str, str]) -> int:
    mapping = {
        "005AAC": palette["primary"],
        "005D7F": palette["secondary"],
        "1DB5CD": palette["accent"],
        "006EE9": palette["primary"],
        "001F3F": palette["text"],
    }
    return sum(
        recolor_shape(shape, mapping)
        for slide in prs.slides
        for shape in slide.shapes
    )


def require_qa_pass(plan: dict[str, Any], base_dir: Path) -> None:
    config = plan.get("vision_review", {})
    if not config.get("enabled") or not config.get("require_pass", True):
        return
    report_path = resolve_path(
        config.get("output", "outputs/image-qa.json"),
        base_dir,
    )
    if report_path is None or not report_path.exists():
        raise FileNotFoundError("vision QA report is required before assembly")
    report = json.loads(report_path.read_text(encoding="utf-8"))
    if not report.get("passed"):
        raise ValueError(f"vision QA did not pass: {report_path}")
    if report.get("reviewer") != "kimi":
        raise ValueError(f"Kimi vision QA report is required: {report_path}")
    expected_ids = {
        slide["id"]
        for section in plan["sections"]
        for slide in section["slides"]
    }
    reviewed_ids = {slide.get("id") for slide in report.get("slides", [])}
    if reviewed_ids != expected_ids:
        raise ValueError(f"vision QA slide set does not match deck plan: {report_path}")


def build_deck(plan_path: Path) -> tuple[Path, Path]:
    plan = load_plan(plan_path)
    errors = validate_plan(plan)
    if errors:
        raise ValueError("\n".join(errors))
    base_dir = plan_path.parent
    require_qa_pass(plan, base_dir)

    template_path = resolve_path(plan.get("template"), base_dir) or DEFAULT_TEMPLATE
    output_path = resolve_path(plan.get("output"), base_dir)
    assert output_path is not None
    output_path.parent.mkdir(parents=True, exist_ok=True)

    minimum_size = configured_minimum_size(plan)
    with tempfile.TemporaryDirectory(prefix="image-avatar-ppt-") as temp_dir:
        selected_template = template_path
        if plan.get("apply_palette_to_template", True):
            selected_template = Path(temp_dir) / "palette-template.pptx"
            prepare_template_with_palette(
                template_path,
                selected_template,
                plan["palette"],
            )

        prs = Presentation(str(selected_template))
        if len(prs.slides) < 5:
            raise ValueError("template must contain at least 5 slides")
        if plan.get("apply_palette_to_template", True):
            recolor_native_shapes(prs, plan["palette"])

        cover_id = prs.slides[0].slide_id
        catalogue_id = prs.slides[1].slide_id
        source_section_slide = prs.slides[2]
        closing_id = prs.slides[4].slide_id
        sections = plan["sections"]
        footer = plan.get("footer")
        ordered_ids: list[int] = []
        inventory: list[dict[str, Any]] = []

        if plan.get("include_cover", True):
            update_cover(prs.slides[0], plan)
            ordered_ids.append(cover_id)
            inventory.append({"type": "cover", "title": plan["cover_title"]})

        if plan.get("include_catalogue", True):
            update_catalogue(prs.slides[1], sections, footer)
            ordered_ids.append(catalogue_id)
            inventory.append({"type": "catalogue", "title": "目录"})

        for section_index, section in enumerate(sections, start=1):
            section_slide = duplicate_slide(prs, source_section_slide)
            replace_matching_text(
                section_slide,
                {
                    "MAIN TITLE": section["main_title"],
                    "文本框 4": section["main_title"],
                    "内容标题": section["subtitle"],
                    "文本框 5": section["subtitle"],
                    "平行数字  交叉现实": footer or "平行数字  交叉现实",
                },
            )
            ordered_ids.append(section_slide.slide_id)
            inventory.append(
                {
                    "type": "chapter",
                    "section": section_index,
                    "main_title": section["main_title"],
                    "subtitle": section["subtitle"],
                }
            )

            for slide in section["slides"]:
                image_path = resolve_path(slide["image"], base_dir)
                assert image_path is not None
                if not image_path.exists():
                    raise FileNotFoundError(f"missing content image: {image_path}")
                width, height = validate_image_geometry(image_path, minimum_size)
                content_slide = add_full_bleed_image_slide(prs, image_path)
                title_render_mode = resolved_title_render_mode(plan, slide)
                if title_render_mode == "native":
                    add_native_content_title(content_slide, plan, slide)
                ordered_ids.append(content_slide.slide_id)
                inventory.append(
                    {
                        "type": "content-image",
                        "section": section_index,
                        "id": slide["id"],
                        "title": slide["title"],
                        "title_render_mode": title_render_mode,
                        "image": str(image_path),
                        "sha256": hashlib.sha256(image_path.read_bytes()).hexdigest(),
                        "dimensions": {"width": width, "height": height},
                    }
                )

        if plan.get("include_closing", True):
            if footer:
                replace_matching_text(
                    prs.slides[4],
                    {"平行数字  交叉现实": footer},
                )
            ordered_ids.append(closing_id)
            inventory.append({"type": "closing", "title": "THANK YOU"})

        keep_ids = set(ordered_ids)
        for slide in list(prs.slides):
            if slide.slide_id not in keep_ids:
                delete_slide_by_id(prs, slide.slide_id)
        reorder_slides(prs, ordered_ids)
        prs.save(str(output_path))

    inventory_path = resolve_path(
        plan.get("inventory_output"),
        base_dir,
    ) or output_path.with_suffix(".inventory.json")
    inventory_path.parent.mkdir(parents=True, exist_ok=True)
    inventory_path.write_text(
        json.dumps(
            {
                "deck": str(output_path),
                "palette": plan["palette"],
                "image_generation": {
                    "provider": plan["image_generation"]["provider"],
                    "model": plan["image_generation"]["model"],
                },
                "slide_count": len(inventory),
                "slides": [
                    {"number": index, **item}
                    for index, item in enumerate(inventory, start=1)
                ],
            },
            ensure_ascii=False,
            indent=2,
        )
        + "\n",
        encoding="utf-8",
    )
    return output_path, inventory_path


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("plan", help="Path to deck-plan JSON")
    args = parser.parse_args()
    output, inventory = build_deck(Path(args.plan).resolve())
    print(f"Created {output}")
    print(f"Created {inventory}")


if __name__ == "__main__":
    main()
