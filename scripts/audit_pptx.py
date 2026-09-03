#!/usr/bin/env python3
"""Statically audit an assembled PPTX for editable copy and font-size contracts."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from deck_utils import (
    iter_slides,
    load_plan,
    prefer_editable_text,
    resolve_path,
    resolved_title_render_mode,
    typography_policy,
)
from validate_plan import validate_plan


def content_slide_indices(plan: dict[str, Any]) -> list[int]:
    """Return zero-based assembled slide indices for plan content slides."""
    index = int(bool(plan.get("include_cover", True)))
    index += int(bool(plan.get("include_catalogue", True)))
    indices: list[int] = []
    for section in plan.get("sections", []):
        index += 1  # chapter divider
        for _slide in section.get("slides", []):
            indices.append(index)
            index += 1
    return indices


def run_sizes(shape: Any) -> list[float]:
    if not getattr(shape, "has_text_frame", False):
        return []
    return [
        float(run.font.size.pt)
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
        if run.text and run.font.size is not None
    ]


def audit_pptx(
    plan: dict[str, Any],
    pptx_path: Path,
) -> dict[str, Any]:
    plan_errors = validate_plan(plan)
    if plan_errors:
        return {
            "passed": False,
            "pptx": str(pptx_path),
            "errors": [f"plan: {error}" for error in plan_errors],
            "slides": [],
        }
    if not pptx_path.exists():
        return {
            "passed": False,
            "pptx": str(pptx_path),
            "errors": [f"PPTX not found: {pptx_path}"],
            "slides": [],
        }

    prs = Presentation(pptx_path)
    typography = typography_policy(plan)
    editable_preferred = prefer_editable_text(plan)
    planned_slides = [slide for _, _, slide in iter_slides(plan)]
    indices = content_slide_indices(plan)
    errors: list[str] = []
    findings: list[dict[str, Any]] = []
    minimum_seen: float | None = None
    editable_shape_count = 0

    if len(indices) != len(planned_slides):
        errors.append("internal content-slide index count does not match the plan")
    if indices and max(indices) >= len(prs.slides):
        errors.append(
            f"assembled slide count {len(prs.slides)} is shorter than the planned sequence"
        )

    for slide_plan, slide_index in zip(planned_slides, indices):
        if slide_index >= len(prs.slides):
            break
        slide = prs.slides[slide_index]
        slide_id = str(slide_plan.get("id", "unknown"))
        finding: dict[str, Any] = {
            "id": slide_id,
            "pptx_slide_number": slide_index + 1,
            "issues": [],
        }
        pictures = [
            shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        ]
        finding["picture_count"] = len(pictures)
        if len(pictures) != 1:
            finding["issues"].append(
                f"expected exactly one full-slide picture, found {len(pictures)}"
            )

        native_titles = [
            shape
            for shape in slide.shapes
            if getattr(shape, "name", "") == "Native Content Title"
        ]
        if resolved_title_render_mode(plan, slide_plan) == "native":
            if len(native_titles) != 1:
                finding["issues"].append(
                    f"expected one native editable title, found {len(native_titles)}"
                )
            elif native_titles[0].text.strip() != str(slide_plan.get("title", "")).strip():
                finding["issues"].append("native title text does not match the plan")
            else:
                sizes = run_sizes(native_titles[0])
                finding["native_title_font_sizes_pt"] = sizes
                if not sizes or min(sizes) < typography["body_font_size_pt"]:
                    finding["issues"].append(
                        "native title is below the regular-text floor of "
                        f"{typography['body_font_size_pt']:g}pt"
                    )

        editable_shapes = [
            shape
            for shape in slide.shapes
            if str(getattr(shape, "name", "")).startswith("Editable Text ")
        ]
        editable_shape_count += len(editable_shapes)
        editable_texts = [shape.text.strip() for shape in editable_shapes]
        finding["editable_text_count"] = len(editable_shapes)
        finding["editable_text"] = editable_texts

        if editable_preferred:
            planned_texts = [str(item) for item in slide_plan.get("exact_text", [])]
            missing = [text for text in planned_texts if editable_texts.count(text) < 1]
            extra = [text for text in editable_texts if text not in planned_texts]
            duplicates = [
                text for text in set(editable_texts) if editable_texts.count(text) > 1
            ]
            if missing:
                finding["issues"].append("missing editable exact text: " + " | ".join(missing))
            if extra:
                finding["issues"].append("unplanned editable text: " + " | ".join(extra))
            if duplicates:
                finding["issues"].append("duplicated editable text: " + " | ".join(duplicates))

        for shape in editable_shapes:
            role = str(getattr(shape, "name", "")).split(":", 1)[-1].strip()
            required = (
                typography["minimum_font_size_pt"]
                if role in {"caption", "label", "note"}
                else typography["body_font_size_pt"]
            )
            sizes = run_sizes(shape)
            if sizes:
                shape_minimum = min(sizes)
                minimum_seen = (
                    shape_minimum
                    if minimum_seen is None
                    else min(minimum_seen, shape_minimum)
                )
            if not sizes or min(sizes) < required:
                finding["issues"].append(
                    f"{shape.name} is below its required {required:g}pt floor"
                )

        if finding["issues"]:
            errors.extend(f"{slide_id}: {issue}" for issue in finding["issues"])
        findings.append(finding)

    expected_slide_count = len(indices) + int(bool(plan.get("include_cover", True)))
    expected_slide_count += int(bool(plan.get("include_catalogue", True)))
    expected_slide_count += len(plan.get("sections", []))
    expected_slide_count += int(bool(plan.get("include_closing", True)))
    if len(prs.slides) != expected_slide_count:
        errors.append(
            f"slide count mismatch: expected {expected_slide_count}, found {len(prs.slides)}"
        )

    return {
        "passed": not errors,
        "pptx": str(pptx_path),
        "errors": errors,
        "metrics": {
            "slide_count": len(prs.slides),
            "content_slide_count": len(planned_slides),
            "editable_text_preferred": editable_preferred,
            "editable_text_shape_count": editable_shape_count,
            "minimum_editable_font_size_pt": minimum_seen,
            "required_minimum_font_size_pt": typography["minimum_font_size_pt"],
            "required_body_font_size_pt": typography["body_font_size_pt"],
        },
        "slides": findings,
    }


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("plan", help="Path to deck-plan JSON")
    parser.add_argument("--pptx", help="PPTX path; defaults to plan.output")
    parser.add_argument("--output", help="Optional JSON report path")
    args = parser.parse_args()

    plan_path = Path(args.plan).resolve()
    plan = load_plan(plan_path)
    pptx_path = resolve_path(args.pptx or plan.get("output"), plan_path.parent)
    if pptx_path is None:
        raise SystemExit("PPTX path is required")
    report = audit_pptx(plan, pptx_path)
    rendered = json.dumps(report, ensure_ascii=False, indent=2) + "\n"
    if args.output:
        output_path = Path(args.output).resolve()
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_text(rendered, encoding="utf-8")
        print(f"Wrote {output_path}; passed={report['passed']}")
    else:
        print(rendered, end="")
    if not report["passed"]:
        raise SystemExit(1)


if __name__ == "__main__":
    main()
